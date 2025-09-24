import os, io, json, zipfile, re
import numpy as np, pandas as pd
import cv2 # opencv-python
import fitz  # PyMuPDF
import pytesseract, ezdxf, openpyxl
from PIL import Image, ImageFilter, ImageEnhance
from docx import Document
from pptx import Presentation
from snownlp import SnowNLP

class FileFullTextExtractor:
  def __init__(self, tesseract_lang='eng+chi_tra'):
    self.lang = tesseract_lang

  def extract_text(self, file_path):
    ext = os.path.splitext(file_path)[1].lower()
    handler = self.get_handler(ext)
    if handler:
      result = handler(file_path)
      if result['status'] == 'ok':
        # 將文字送到 clean_text 處理
        # result['text'] = self.clean_text(result['text'])
        result['text'] = result['text']
      return result
    return {'status': 'fail', 'error': f"Unsupported file type: {ext}"}

  def get_handler(self, ext):
    return {
      ".pdf": self.extract_from_pdf,
      ".docx": self.extract_from_docx,
      ".doc": self.extract_from_docx,
      ".xlsx": self.extract_from_excel,
      ".xls": self.extract_from_excel,
      ".pptx": self.extract_from_pptx,
      ".ppt": self.extract_from_pptx,
      ".jpg": self.extract_from_image,
      ".jpeg": self.extract_from_image,
      ".png": self.extract_from_image,
      ".bmp": self.extract_from_image,
      ".tif": self.extract_from_image,
      ".tiff": self.extract_from_image,
      ".txt": self.extract_from_txt,
      ".csv": self.extract_from_csv,
      ".json": self.extract_from_json,
      #".dwg": self.extract_from_dwg_or_dxf,
      ".dxf": self.extract_from_dwg_or_dxf,
      ".zip": self.extract_from_zip,
    }.get(ext)

  def extract_from_pdf(self, file_path):
    all_text = ''

    # 1️⃣ 嘗試文字型 PDF（含表格座標排序）
    try:
      doc = fitz.open(file_path)
    except Exception as e:
      return {'status': 'fail', 'error': f'Cannot open PDF: {e}'}

    # 2️⃣ OCR 提取
    for page_num in range(len(doc)):
      page = doc.load_page(page_num)
      page_text = ""

      # 1️⃣ 嘗試直接取文字
      try:
        page_text = page.get_text("text")
        if page_text.strip():
          all_text += f"\n{page_text}"
      except Exception as e:
        print( f"⚠ fitz get_text error page {page_num+1}: {e}")
      # 2️⃣ OCR fallback
      if not page_text.strip():
        try:
          pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))  # 放大 3倍
          img = Image.open(io.BytesIO(pix.tobytes("png")))
          ocr_text = pytesseract.image_to_string(img, lang=self.lang, config='--psm 6')
          if ocr_text.strip():
            all_text += f"\n{ocr_text}"
        except Exception as e:
          print(f"⚠ OCR error page {page_num+1}: {e}")

    # 回傳結果
    all_text = all_text.strip()
    if all_text:
      return {'status': 'ok', 'text': self.clean_ocr_text(all_text)}
    else:
      return {'status': 'fail', 'error': 'PDF無法提取文字'}


  def extract_from_docx(self, file_path):
    try:
      doc = Document(file_path)
      all_text = ""

      # 1️⃣ 段落文字
      for p in doc.paragraphs:
        if p.text.strip():
          all_text += p.text.strip() + "\n"

      # 2️⃣ 表格文字
      for table in doc.tables:
        for row in table.rows:
          row_text = "\n".join(cell.text.strip() for cell in row.cells if cell.text.strip())
          if row_text:
            all_text += row_text + "\n"

      # 3️⃣ 內嵌圖片 OCR
      for rel in doc.part._rels:
        rel_obj = doc.part._rels[rel]
        if "image" in rel_obj.target_ref:
          image_data = rel_obj.target_part.blob
          img = Image.open(io.BytesIO(image_data))
          ocr_result = self.extract_from_image(img)
          if ocr_result['status'] == 'ok' and ocr_result['text'].strip():
            all_text += "\n" + ocr_result['text'].strip() + "\n"

      # 4️⃣ 清理文字
      return {'status': 'ok', 'text': self.clean_text(all_text.strip())}

    except Exception as e:
      return {'status': 'fail', 'error': f'DOCX extract error: {e}'}


  def extract_from_excel(self, file_path):
    try:
      # 1️⃣ 讀取所有 sheet
      df_list = pd.read_excel(file_path, sheet_name=None, dtype=str)
      all_text = ''

      for sheet, df in df_list.items():
        # 將表格文字加入
        all_text += df.fillna('').to_string(index=False) + "\n"

      # 2️⃣ 處理內嵌圖片
      wb = openpyxl.load_workbook(file_path)
      for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for image in ws._images:  # openpyxl 內嵌圖片物件
          img = image._data() if hasattr(image, "_data") else image.ref
          # image 是 PIL.Image 可以直接 OCR，如果不是則用 BytesIO 轉
          if isinstance(img, bytes):
            img = Image.open(io.BytesIO(img))
          ocr_result = self.extract_from_image(img)
          if ocr_result['status'] == 'ok' and ocr_result['text'].strip():
            all_text += "\n" + ocr_result['text'].strip() + "\n"

      # 3️⃣ 清理文字
      return {'status': 'ok', 'text': self.clean_text(all_text.strip())}

    except Exception as e:
      return {'status': 'fail', 'error': f'Excel parse error: {e}'}


  def extract_from_pptx(self, file_path):
    try:
      prs = Presentation(file_path)
      all_text = ''

      for slide in prs.slides:
        for shape in slide.shapes:
          # 文字
          if hasattr(shape, "text") and shape.text.strip():
            all_text += shape.text.strip() + "\n"

          # 圖片
          if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE == 13
            image = shape.image
            img =Image.open(io.BytesIO(image.blob))
            ocr_result = self.extract_from_image(img)
            if ocr_result.get('status') == 'ok' and ocr_result.get('text').strip():
              all_text += ocr_result['text'].strip() + "\n"

      return {'status': 'ok', 'text': self.clean_text(all_text.strip())}

    except Exception as e:
      return {'status': 'fail', 'error': f'PowerPoint parse error: {e}'}

  @staticmethod
  def preprocess_image(file_path):
    img = Image.open(file_path)
    img = img.convert('L') # 灰階

    # 放大
    width, height = img.size
    img = img.resize((width*2, height*2), Image.Resampling.LANCZOS)

    # 提升對比
    enhancer = ImageEnhance.Contrast(img)
    img = enhancer.enhance(2.0)

    # 銳化
    img = img.filter(ImageFilter.SHARPEN)

    # 轉成 numpy array 做二值化
    img_cv = np.array(img)
    _, img_bin = cv2.threshold(img_cv, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    return Image.fromarray(img_bin)

  def extract_from_image(self, file_input):
    """
    file_input: 可以是檔案路徑(str) 或 PIL.Image.Image
    """
    try:
      # 如果是路徑就先 preprocess
      if isinstance(file_input, str):
          img = self.preprocess_image(file_input)
      else:  # 假設已經是 PIL.Image
          img = file_input

      # 嘗試多種 PSM
      text = pytesseract.image_to_string(img, lang=self.lang, config='--psm 6')
      return {'status': 'ok', 'text': self.clean_ocr_text(text)}
    except Exception as e:
      return {'status': 'fail', 'error': f'OCR error in image: {e}'}

  def extract_from_txt(self, file_path):
    try:
      with open(file_path, encoding='utf-8') as f:
        text = f.read()
      return {'status': 'ok', 'text': self.clean_text(text)}
    except Exception as e:
      return {'status': 'fail', 'error': f'TXT parse error: {e}'}

  def extract_from_csv(self, file_path):
    try:
      text = pd.read_csv(file_path).to_string(index=False)
      return {'status': 'ok', 'text': self.clean_text(text)}
    except Exception as e:
      return {'status': 'fail', 'error': f'CSV parse error: {e}'}

  def extract_from_json(self, file_path):
    try:
      with open(file_path, encoding='utf-8') as f:
        text = json.dumps(json.load(f), indent=2, ensure_ascii=False)
      return {'status': 'ok', 'text': self.clean_text(text)}
    except Exception as e:
      return {'status': 'fail', 'error': f'JSON parse error: {e}'}

  def extract_from_dwg_or_dxf(self, file_path):
    """
    支援 DXF，自動檢測 DWG/DXF：
    - DXF: 提取文字
    - DWG: 提示需轉 DXF
    """
    try:
      ext = os.path.splitext(file_path)[1].lower()

      if ext == '.dxf':
        doc = ezdxf.readfile(file_path)
        msp = doc.modelspace()
        texts = [e.dxf.text for e in msp if e.dxftype() in ['TEXT', 'MTEXT']]
        return {'status': 'ok', 'text': self.clean_text('\n'.join(texts))}

      elif ext == '.dwg':
        return {'status': 'fail', 'error': f'DWG file detected: {file_path}. Please convert to DXF first.'}

      else:
        return {'status': 'fail', 'error': f'Unsupported CAD file type: {ext}'}

    except Exception as e:
      return {'status': 'fail', 'error': f'CAD parse error: {e}'}

  def extract_from_zip(self, file_path):
    try:
      result_text = ''
      with zipfile.ZipFile(file_path, 'r') as zip_ref:
        for file_name in zip_ref.namelist():
          # result_text += f'--- ZIP ENTRY: {file_name} ---\n'
          ext = os.path.splitext(file_name)[1].lower()
          handler = self.get_handler(ext)
          if handler:
            # 讀取檔案內容到暫存檔
            with zip_ref.open(file_name) as f:
              tmp_path = f'./temp_{os.path.basename(file_name)}'
              with open(tmp_path, 'wb') as tmp_f:
                tmp_f.write(f.read())
            # 呼叫對應 extract 函式
            extracted = handler(tmp_path)
            if extracted['status'] == 'ok':
              result_text += extracted['text'] + '\n'
            # else:
            #   result_text += f"[Failed to extract {file_name}]: {extracted.get('error', '')}\n"
            os.remove(tmp_path)
          # else:
            result_text += f"[No handler for {file_name}]\n"
      return {'status': 'ok', 'text': self.clean_text(result_text)}
    except Exception as e:
      return {'status': 'fail', 'error': f'ZIP parse error: {e}'}

  def clean_text(self, text):
    # 移除頁面標題及多餘符號，並去掉所有空白與換行
    lines = text.splitlines()
    cleaned_lines = []
    for line in lines:
      line = line.strip()
      if not line:
        continue

      # 移除頁面標題
      if re.match(r'^--- Page \d+ - Rotate \d+° ---$', line):
        continue
      if re.match(r'^--- Sheet: Table \d+ ---$', line):
        continue
      if re.match(r'^--- Page \d+ Text ---$', line):
        continue

      # 移除只包含符號或亂碼的行
      if re.match(r'^[\W\dA-Za-z]{0,5}$', line):
        continue
      # 移除超長亂碼行
      if len(re.sub(r'\s', '', line)) > 100 and re.match(r'^[\W\da-zA-Z]+$', line):
        continue
      # 移除所有空白
      line = re.sub(r'\s+', '', line)
      # 去掉多餘的符號
      line = re.sub(r'[^0-9\u4e00-\u9fffA-Za-z/-]', '', line)
      cleaned_lines.append(line)

    return ''.join(cleaned_lines)

  def clean_ocr_text(self, text, min_len=6):
    """
    將長串 OCR 文本斷句、過濾噪音，輸出乾淨文字
    min_len: 最短字數，過濾太短或無意義片段
    """
    if not text:
      return ''

    # 移除連續空白
    text = re.sub(r' +', '', text)

    # 用 SnowNLP 自動斷句
    s = SnowNLP(text)
    sentences = s.sentences

    # 過濾短句、全為符號或單字的片段
    clean_sentences = []
    for sent in sentences:
      sent = sent.strip()
      if len(sent) >= min_len and not re.fullmatch(r'[\d\W_]+', sent):
        clean_sentences.append(sent)
    # 移除過長的連續數字或英文編號一
    clean_sentences = list(map(lambda s: re.sub(r'[A-Za-z0-9一\/-]{20,}', '', s), sentences))

    # 合併成單一文字
    clean_text = self.clean_text(''.join(clean_sentences))
    return clean_text if clean_text else text.strip()
